from pathlib import Path
from pptx import Presentation

from app.parsers.base import BaseParser


class PPTXParser(BaseParser):
    """PPTX presentation parser using python-pptx."""

    def supports(self, file_type: str) -> bool:
        return file_type.lower() == "pptx"

    def parse(self, file_path: Path) -> str:
        """Extract text from PPTX presentation."""
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            if len(slide_text) > 1:  # More than just the slide marker
                text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts)

    def get_metadata(self, file_path: Path) -> dict:
        """Extract PPTX metadata."""
        metadata = super().get_metadata(file_path)

        try:
            prs = Presentation(file_path)
            metadata["slide_count"] = len(prs.slides)

            core_props = prs.core_properties
            if core_props.title:
                metadata["title"] = core_props.title
            if core_props.author:
                metadata["author"] = core_props.author
        except Exception:
            pass

        return metadata
